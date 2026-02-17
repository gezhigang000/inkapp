#!/usr/bin/env python3
"""
Sidecar 入口：接收 Tauri 前端命令，输出 JSON Lines 进度。
用法: echo '{"action":"generate","mode":"daily"}' | python3 sidecar_main.py
"""
import sys
import os
import io

# Windows 下强制 UTF-8 编码，避免中文乱码
if sys.platform == "win32":
    # stdin: Rust 发送 UTF-8 字节，Python 默认按 GBK 解码会乱码
    sys.stdin = io.TextIOWrapper(sys.stdin.buffer, encoding="utf-8", errors="replace")
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8", errors="replace")
    os.environ.setdefault("PYTHONIOENCODING", "utf-8")

import json
import logging
from datetime import datetime

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))

# PyInstaller 打包后，数据文件在 sys._MEIPASS 目录下
if getattr(sys, 'frozen', False) and hasattr(sys, '_MEIPASS'):
    PROJECT_ROOT = sys._MEIPASS
else:
    PROJECT_ROOT = os.path.dirname(SCRIPT_DIR)

# 确保 scripts 目录在 import 路径中
if SCRIPT_DIR not in sys.path:
    sys.path.insert(0, SCRIPT_DIR)

# Pillow 等依赖可能在 pylib 目录
PYLIB_DIR = os.path.join(PROJECT_ROOT, "pylib")
if os.path.exists(PYLIB_DIR) and PYLIB_DIR not in sys.path:
    sys.path.insert(0, PYLIB_DIR)

# ============================================================
# 本地日志 & 缓存目录
# 跨平台：Windows 用 %APPDATA%/Ink，macOS/Linux 用 ~/.ink
# ============================================================
if sys.platform == "win32":
    _app_data = os.environ.get("APPDATA", os.path.expanduser("~"))
    INK_HOME = os.path.join(_app_data, "Ink")
else:
    INK_HOME = os.path.join(os.path.expanduser("~"), ".ink")
LOG_DIR = os.path.join(INK_HOME, "logs")
CACHE_DIR = os.path.join(INK_HOME, "cache")
os.makedirs(LOG_DIR, exist_ok=True)
os.makedirs(CACHE_DIR, exist_ok=True)

# 按天滚动的日志文件
_log_file = os.path.join(LOG_DIR, f"{datetime.now().strftime('%Y-%m-%d')}.log")
logging.basicConfig(
    filename=_log_file,
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    datefmt="%H:%M:%S",
)
logger = logging.getLogger("ink")


def emit(event_type, **kwargs):
    """输出一行 JSON 事件到 stdout，同时写入日志"""
    event = {"type": event_type, **kwargs}
    print(json.dumps(event, ensure_ascii=False), flush=True)
    # 写入本地日志
    msg = kwargs.get("message", "")
    if event_type == "error":
        logger.error("emit %s: %s %s", event_type, kwargs.get("code", ""), msg)
    elif event_type == "result":
        logger.info("emit %s: status=%s %s", event_type, kwargs.get("status", ""), msg)
    else:
        logger.info("emit %s: %s", event_type, msg)


def _get_cover_kwargs(params):
    """从 payload 提取封面设置参数"""
    return {
        "color_style": params.get("cover_color_style", "random"),
        "pattern_style": params.get("cover_pattern_style", "random"),
        "show_title": params.get("cover_show_title", True),
        "subtitle": params.get("cover_subtitle", "Ink"),
    }


def _handle_translate_inplace(params, config, file_formats, topic, timestamp,
                              output_dir, header_html, footer_html):
    """原格式翻译：保持文档格式和样式，只替换文字内容。"""
    import re as re_mod
    from translate_inplace import translate_file_inplace
    from daily_ai_news import (
        generate_cover_image, save_article, make_timestamp,
        pick_daily_variation, append_footer,
    )

    # 从 topic 提取目标语言（如 "翻译为英文" → "英文"）
    lang_match = re_mod.search(
        r'(?:翻译[为成]?|译[为成]?|translate\s*(?:to|into)?)\s*(.+)',
        topic, re_mod.IGNORECASE)
    target_lang = lang_match.group(1) if lang_match else "英文"
    logger.info("Inplace translation: target_lang=%s, files=%d",
                target_lang, len(file_formats))

    article_dir = os.path.join(output_dir, timestamp)
    os.makedirs(article_dir, exist_ok=True)

    translated_files = []
    all_preview_texts = []

    for finfo in file_formats:
        src_path = finfo.get("path", "")
        fname = finfo.get("name", "")
        ext = finfo.get("ext", "").lower()

        if not src_path or not os.path.exists(src_path):
            logger.warning("File not found: %s", src_path)
            continue
        if ext not in ("docx", "doc", "pptx", "ppt", "pdf"):
            continue

        # 输出文件名：原文件名加 _translated 后缀
        base, dot_ext = os.path.splitext(fname)
        out_name = f"{base}_translated{dot_ext}"
        out_path = os.path.join(article_dir, out_name)

        emit("progress", stage="translating",
             message=f"正在翻译 {fname}...", percent=20)

        result_path = translate_file_inplace(
            src_path, out_path, target_lang, config, emit_fn=emit)

        if result_path:
            translated_files.append(out_name)
            logger.info("Translated: %s → %s", fname, out_name)
            # 收集预览文本
            try:
                preview = _extract_preview_text(result_path, ext)
                if preview:
                    all_preview_texts.append(
                        f"<h3>{fname} → {target_lang}</h3>\n{preview}")
            except Exception:
                pass
        else:
            emit("progress", stage="translating",
                 message=f"翻译 {fname} 失败，回退到文本翻译模式")
            logger.warning("Inplace translation failed for %s", fname)

    if not translated_files:
        emit("error", code="TRANSLATE_FAILED",
             message="翻译失败，无法处理上传的文件")
        return

    # 生成 HTML 预览页面
    emit("progress", stage="saving", message="正在保存...", percent=80)
    preview_html = _build_preview_html(all_preview_texts, target_lang)
    if header_html:
        preview_html = header_html + preview_html
    if footer_html:
        preview_html = preview_html + footer_html

    html_path = os.path.join(article_dir, f"{timestamp}-report.html")
    with open(html_path, "w", encoding="utf-8") as f:
        f.write(preview_html)

    # 生成封面图
    emit("progress", stage="cover", message="正在生成封面图...", percent=90)
    today = datetime.now().strftime("%Y-%m-%d")
    variation = pick_daily_variation(today)
    title = f"文档翻译 - {target_lang}"
    cover_kwargs = _get_cover_kwargs(params)
    cover_path = generate_cover_image(
        timestamp, title, topic, article_dir,
        cover_theme=variation.get("cover_theme"),
        **cover_kwargs,
    )

    # 保存元数据
    converted_path = os.path.join(article_dir, translated_files[0]) \
        if translated_files else ""
    file_type = os.path.splitext(translated_files[0])[1].lstrip(".") \
        if translated_files else "html"
    metadata = {
        "title": title,
        "date": timestamp[:8],
        "mode": "topic",
        "topic": topic,
        "status": "generated",
        "provider": config.get("LLM_PROVIDER", "deepseek"),
        "file_type": file_type,
        "converted_path": converted_path,
        "output_files": translated_files,
        "articles": [
            {"title": title, "path": html_path,
             "cover": str(cover_path) if cover_path else ""}
        ],
    }
    meta_path = os.path.join(article_dir, f"{timestamp}-metadata.json")
    with open(meta_path, "w", encoding="utf-8") as f:
        json.dump(metadata, f, ensure_ascii=False, indent=2)

    emit("progress", stage="done", message="翻译完成！", percent=100)
    emit("result", status="success", title=title,
         article_path=html_path, metadata_path=meta_path,
         file_type=file_type, article_count=1)


def _extract_preview_text(file_path, ext):
    """从翻译后的文件中提取预览文本（前几段）"""
    if ext in ("docx", "doc"):
        from docx import Document
        doc = Document(file_path)
        lines = [p.text for p in doc.paragraphs[:20] if p.text.strip()]
        return "<br>".join(lines[:10])
    elif ext in ("pptx", "ppt"):
        from pptx import Presentation
        prs = Presentation(file_path)
        lines = []
        for slide in prs.slides[:5]:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            lines.append(para.text)
        return "<br>".join(lines[:10])
    elif ext == "pdf":
        try:
            import fitz
            doc = fitz.open(file_path)
            lines = []
            for page in doc[:3]:
                text = page.get_text()
                if text:
                    lines.extend(text.strip().split("\n")[:5])
            doc.close()
            return "<br>".join(lines[:10])
        except Exception:
            return ""
    return ""


def _build_preview_html(preview_texts, target_lang):
    """构建翻译预览 HTML"""
    import html as html_mod
    font = "-apple-system, BlinkMacSystemFont, 'PingFang SC', sans-serif"
    parts = [
        f'<section style="padding:20px 0;font-family:{font};">',
        f'<h2 style="font-size:18px;color:#333;margin-bottom:16px;">'
        f'翻译预览（{html_mod.escape(target_lang)}）</h2>',
        '<p style="font-size:13px;color:#999;margin-bottom:20px;">'
        '以下为翻译后文档的部分内容预览，完整文件请点击「打开文件夹」查看。</p>',
    ]
    for text in preview_texts:
        escaped = html_mod.escape(text)
        parts.append(
            f'<div style="padding:16px;background:#f8f9fa;border-radius:8px;'
            f'margin-bottom:12px;font-size:14px;line-height:1.8;color:#333;">'
            f'{escaped}</div>'
        )
    parts.append('</section>')
    return "\n".join(parts)


def handle_generate(params):
    """处理文章生成请求，集成 daily_ai_news.py 核心逻辑"""
    from datetime import datetime
    from daily_ai_news import (
        generate_article, generate_cover_image,
        extract_html, extract_title, split_article_if_needed,
        make_timestamp, run_video_analysis, save_article,
        pick_daily_variation, append_footer,
    )

    logger.info("=== generate start === mode=%s topic=%s provider=%s",
                params.get("mode"), params.get("topic", "")[:60],
                params.get("provider", "?"))
    emit("progress", stage="init", message="正在加载配置...")

    # 配置完全从前端参数构建，不依赖 config.env 文件
    config = {}
    if params.get("provider"):
        config["LLM_PROVIDER"] = params["provider"]
    for key in ["DEEPSEEK_API_KEY", "GLM_API_KEY", "DOUBAO_API_KEY",
                "KIMI_API_KEY", "OPENAI_API_KEY",
                "DEEPSEEK_MODEL", "GLM_MODEL", "DOUBAO_MODEL",
                "KIMI_MODEL", "OPENAI_MODEL",
                "TAVILY_API_KEY", "SERPAPI_API_KEY",
                "SEARCH_PROVIDER", "OUTPUT_DIR"]:
        if params.get(key):
            config[key] = params[key]

    mode = params.get("mode", "daily")
    topic = params.get("topic", "")
    video_url = params.get("video_url", "")
    file_contents = params.get("file_contents", "")  # 上传文件提取的文本
    template_prompt = params.get("template_prompt", "")  # 模板自定义提示词
    file_formats = params.get("file_formats", None)  # 上传文件格式信息
    header_html = params.get("header_html", "")  # 文章头部模板
    footer_html = params.get("footer_html", "")  # 文章尾部模板
    # OSS 云存储配置
    oss_config = {}
    for k in ["oss_bucket", "oss_endpoint", "oss_access_key_id", "oss_access_key_secret"]:
        if params.get(k):
            oss_config[k] = params[k]
    default_output = os.path.join(INK_HOME, "articles")
    output_dir = config.get("OUTPUT_DIR", default_output)
    timestamp = make_timestamp()

    try:  # noqa: E501 — 捕获 SystemExit（daily_ai_news 内部 sys.exit）
        # ---------- 视频分析模式 ----------
        if mode == "video" and video_url:
            emit("progress", stage="analyzing", message="正在分析视频...")
            args = {"local": True, "publish": False, "video": video_url, "topic": None}
            run_video_analysis(video_url, config, args)
            emit("progress", stage="done", message="视频分析完成！", percent=100)
            emit("result", status="success", title="视频分析完成", article_path=output_dir)
            return

        # ---------- 原格式翻译模式 ----------
        # 检测：翻译模板 + 有上传文件（docx/pptx/pdf）
        template_id = params.get("template_id", "")
        is_translate = template_id == "translate"
        has_translatable_file = (
            is_translate and file_formats and
            any(f.get("ext", "").lower() in ("docx", "doc", "pptx", "ppt", "pdf")
                for f in file_formats)
        )
        if has_translatable_file:
            _handle_translate_inplace(
                params, config, file_formats, topic, timestamp,
                output_dir, header_html, footer_html)
            return

        # ---------- 文章生成模式（daily / topic） ----------
        today = datetime.now().strftime("%Y-%m-%d")
        variation = pick_daily_variation(today)
        effective_topic = topic if topic else variation.get("topic")

        # 如果有上传文件内容，保持 topic 干净，file_contents 单独传递
        if file_contents and not effective_topic:
            effective_topic = "数据分析报告"

        emit("progress", stage="generating", message="正在生成文章...", percent=20)
        html_content = generate_article(topic=effective_topic, config=config,
                                        custom_prompt=template_prompt,
                                        file_contents=file_contents)
        if not html_content:
            emit("error", code="GENERATION_FAILED", message="文章生成失败，未获得输出")
            return

        # 缓存 LLM 原始输出，方便排查和复用
        cache_file = os.path.join(CACHE_DIR, f"{timestamp}-raw.html")
        try:
            with open(cache_file, "w", encoding="utf-8") as cf:
                cf.write(html_content)
            logger.info("LLM raw output cached: %s (%d chars)", cache_file, len(html_content))
        except OSError:
            pass

        emit("progress", stage="processing", message="正在处理文章...", percent=50)

        # 提取 HTML
        extracted = extract_html(html_content)
        if extracted:
            html_content = extracted

        # 提取标题
        title = extract_title(html_content) or f"AI 资讯 {timestamp}"

        # 拆分长文
        articles = split_article_if_needed(html_content, title)
        is_series = len(articles) > 1

        emit("progress", stage="saving", message="正在保存文章...", percent=60)

        filepaths = []
        img_paths = []
        for idx, (part_title, part_html) in enumerate(articles):
            part_html = append_footer(part_html)
            # 插入用户自定义头部/尾部模板
            if header_html:
                part_html = header_html + part_html
            if footer_html:
                part_html = part_html + footer_html
            suffix = f"-part{idx + 1}" if is_series else ""
            filepath = save_article(timestamp, part_html, output_dir, suffix=suffix)
            filepaths.append(str(filepath))

            emit("progress", stage="cover", message=f"正在生成封面图...", percent=70 + idx * 10)
            cover_kwargs = _get_cover_kwargs(params)
            cover_path = generate_cover_image(
                f"{timestamp}{suffix}", part_title,
                effective_topic or "", output_dir,
                cover_theme=variation.get("cover_theme"),
                **cover_kwargs,
            )
            img_paths.append(str(cover_path) if cover_path else "")

        # 同步到 OSS
        if len(oss_config) == 4:
            emit("progress", stage="uploading", message="正在同步到云端...", percent=90)
            try:
                for i in range(len(articles)):
                    # 上传文章 HTML
                    oss_article_key = f"articles/{timestamp}/{os.path.basename(filepaths[i])}"
                    _upload_to_oss(filepaths[i], oss_article_key, oss_config)
                    # 上传封面图
                    if img_paths[i] and os.path.exists(img_paths[i]):
                        oss_cover_key = f"articles/{timestamp}/{os.path.basename(img_paths[i])}"
                        _upload_to_oss(img_paths[i], oss_cover_key, oss_config)
                emit("progress", stage="uploading", message="云端同步完成", percent=95)
            except Exception as e:
                emit("progress", stage="uploading",
                     message=f"OSS 同步失败（不影响本地文件）: {e}")

        # 格式转换：如果有上传文件，生成与上传文件相同格式的输出
        file_type = "html"
        converted_files = []
        converted_path = ""
        if file_formats:
            target_ext = file_formats[0].get("ext", "").lower()
            if target_ext and target_ext not in ("txt", "md", "csv"):
                article_dir = os.path.dirname(filepaths[0]) if filepaths else output_dir
                emit("progress", stage="converting",
                     message=f"正在转换为 .{target_ext} 格式...", percent=85)
                converted = _convert_html_to_format(
                    html_content, target_ext, article_dir, timestamp)
                if converted:
                    file_type = target_ext
                    converted_path = converted
                    converted_files.append(os.path.basename(converted))
                    logger.info("Converted to %s: %s", target_ext, converted)

        # 保存元数据
        meta_dir = os.path.dirname(filepaths[0]) if filepaths else output_dir
        metadata = {
            "title": title,
            "date": timestamp[:8],  # YYYYMMDD
            "mode": mode,
            "topic": effective_topic or "",
            "status": "generated",
            "provider": config.get("LLM_PROVIDER", "claude"),
            "file_type": file_type,
            "converted_path": converted_path,
            "output_files": converted_files,
            "articles": [
                {"title": articles[i][0], "path": filepaths[i], "cover": img_paths[i]}
                for i in range(len(articles))
            ],
        }
        meta_path = os.path.join(meta_dir, f"{timestamp}-metadata.json")
        with open(meta_path, "w", encoding="utf-8") as f:
            json.dump(metadata, f, ensure_ascii=False, indent=2)

        emit("progress", stage="done", message="生成完成！", percent=100)
        emit("result", status="success", title=title,
             article_path=filepaths[0], metadata_path=meta_path,
             file_type=file_type, article_count=len(articles))

    except SystemExit as e:
        logger.error("generate SystemExit code=%s", e.code)
        emit("error", code="GENERATION_ERROR", message=f"生成过程异常退出 (code={e.code})")
    except Exception as e:
        logger.exception("generate exception")
        emit("error", code="GENERATION_ERROR", message=str(e))


def handle_agent_generate(params):
    """处理 Agent 模式生成请求：多轮工具调用"""
    import shutil
    from agent_loop import run_agent_loop, init_workspace
    from daily_ai_news import (
        extract_html, extract_title, append_footer,
        make_timestamp, save_article, generate_cover_image,
        pick_daily_variation,
    )

    logger.info("=== agent_generate start === topic=%s provider=%s",
                params.get("topic", "")[:60], params.get("provider", "?"))
    emit("progress", stage="init", message="正在初始化 Agent...")

    # 构建 config（同 handle_generate）
    config = {}
    if params.get("provider"):
        config["LLM_PROVIDER"] = params["provider"]
    for key in ["DEEPSEEK_API_KEY", "GLM_API_KEY", "DOUBAO_API_KEY",
                "KIMI_API_KEY", "OPENAI_API_KEY",
                "DEEPSEEK_MODEL", "GLM_MODEL", "DOUBAO_MODEL",
                "KIMI_MODEL", "OPENAI_MODEL",
                "TAVILY_API_KEY", "SERPAPI_API_KEY",
                "SEARCH_PROVIDER", "OUTPUT_DIR"]:
        if params.get(key):
            config[key] = params[key]

    topic = params.get("topic", "")
    file_contents = params.get("file_contents", "")
    template_prompt = params.get("template_prompt", "")
    header_html = params.get("header_html", "")
    footer_html = params.get("footer_html", "")
    file_formats = params.get("file_formats", None)

    # OSS 配置
    oss_config = {}
    for k in ["oss_bucket", "oss_endpoint", "oss_access_key_id", "oss_access_key_secret"]:
        if params.get(k):
            oss_config[k] = params[k]

    default_output = os.path.join(INK_HOME, "articles")
    output_dir = config.get("OUTPUT_DIR", default_output)
    timestamp = make_timestamp()

    try:
        # 初始化 workspace
        workspace = init_workspace(timestamp)
        emit("progress", stage="agent", message=f"工作区: {workspace}")

        # 写入上传文件文本
        if file_contents:
            input_path = os.path.join(workspace, "input", "uploaded_data.txt")
            with open(input_path, "w", encoding="utf-8") as f:
                f.write(file_contents)

        # 复制原始文件到 workspace/input/（翻译场景需要原始二进制）
        if file_formats:
            for finfo in file_formats:
                src_path = finfo.get("path", "")
                if src_path and os.path.exists(src_path):
                    try:
                        dst_path = os.path.join(workspace, "input", finfo["name"])
                        shutil.copy2(src_path, dst_path)
                        logger.info("Copied file to workspace: %s", finfo["name"])
                    except (IOError, OSError) as e:
                        logger.warning("Failed to copy %s: %s", finfo["name"], e)

        if not topic:
            topic = "深度调研报告"

        # 运行 Agent 循环
        # 轮次由模板定义，前端透传 max_turns
        turns = params.get("max_turns", 15)
        html_content = run_agent_loop(
            topic=topic,
            config=config,
            emit_fn=emit,
            workspace=workspace,
            template_prompt=template_prompt,
            file_contents=file_contents,
            file_formats=file_formats,
            max_turns=turns,
        )

        if not html_content:
            emit("error", code="AGENT_FAILED", message="Agent 未生成输出内容")
            return

        # 后处理：检查 workspace/output/ 下是否有原格式文件（翻译场景）
        output_files = []
        ws_output = os.path.join(workspace, "output")
        if os.path.exists(ws_output):
            for fname in os.listdir(ws_output):
                if fname != "article.html" and not fname.endswith(".tmp"):
                    output_files.append(fname)

        # 如果 Agent 没有生成原格式文件，sidecar 自动从 HTML 转换
        if file_formats and not output_files:
            target_ext = file_formats[0].get("ext", "").lower()
            if target_ext and html_content:
                emit("progress", stage="converting",
                     message=f"正在转换为 .{target_ext} 格式...")
                converted = _convert_html_to_format(
                    html_content, target_ext, ws_output, timestamp)
                if converted:
                    output_files.append(os.path.basename(converted))
                    logger.info("Auto-converted HTML to %s: %s",
                                target_ext, converted)

        # 缓存 Agent 原始输出
        cache_file = os.path.join(CACHE_DIR, f"{timestamp}-agent-raw.html")
        try:
            with open(cache_file, "w", encoding="utf-8") as cf:
                cf.write(html_content)
        except OSError:
            pass

        emit("progress", stage="processing", message="正在处理文章...", percent=50)

        # 提取 HTML
        extracted = extract_html(html_content)
        if extracted:
            html_content = extracted

        title = extract_title(html_content) or f"Agent 创作 {timestamp}"

        emit("progress", stage="saving", message="正在保存文章...", percent=60)

        # 添加页脚
        html_content = append_footer(html_content)
        if header_html:
            html_content = header_html + html_content
        if footer_html:
            html_content = html_content + footer_html

        filepath = save_article(timestamp, html_content, output_dir)

        emit("progress", stage="cover", message="正在生成封面图...", percent=70)
        today = datetime.now().strftime("%Y-%m-%d")
        variation = pick_daily_variation(today)
        cover_kwargs = _get_cover_kwargs(params)
        cover_path = generate_cover_image(
            timestamp, title, topic, output_dir,
            cover_theme=variation.get("cover_theme"),
            **cover_kwargs,
        )
        article_dir = os.path.dirname(str(filepath))
        file_type = "html"
        for fname in output_files:
            src = os.path.join(ws_output, fname)
            dst = os.path.join(article_dir, fname)
            shutil.copy2(src, dst)
            ext = os.path.splitext(fname)[1].lower()
            if ext in (".docx", ".xlsx", ".pdf"):
                file_type = ext[1:]  # docx/xlsx/pdf
            logger.info("Copied output file: %s", fname)

        # 保存元数据
        metadata = {
            "title": title,
            "date": timestamp[:8],
            "mode": "agent",
            "topic": topic,
            "status": "generated",
            "provider": config.get("LLM_PROVIDER", "deepseek"),
            "file_type": file_type,
            "output_files": output_files,
            "articles": [
                {"title": title, "path": str(filepath),
                 "cover": str(cover_path) if cover_path else ""}
            ],
        }
        meta_path = os.path.join(article_dir, f"{timestamp}-metadata.json")
        with open(meta_path, "w", encoding="utf-8") as f:
            json.dump(metadata, f, ensure_ascii=False, indent=2)

        emit("progress", stage="done", message="Agent 创作完成！", percent=100)
        emit("result", status="success", title=title,
             article_path=str(filepath), metadata_path=meta_path,
             file_type=file_type, article_count=1)

    except SystemExit as e:
        logger.error("agent_generate SystemExit code=%s", e.code)
        emit("error", code="AGENT_ERROR", message=f"Agent 异常退出 (code={e.code})")
    except Exception as e:
        logger.exception("agent_generate exception")
        emit("error", code="AGENT_ERROR", message=str(e))


def handle_validate_key(params):
    """验证 API Key 有效性：向对应平台发送轻量请求"""
    provider = params.get("provider", "")
    api_key = params.get("api_key", "")

    if not provider or not api_key:
        emit("error", code="MISSING_PARAMS", message="缺少 provider 或 api_key")
        return

    import requests

    endpoints = {
        "deepseek": "https://api.deepseek.com/v1/models",
        "glm": "https://open.bigmodel.cn/api/paas/v4/models",
        "doubao": "https://ark.cn-beijing.volces.com/api/v3/models",
        "kimi": "https://api.moonshot.cn/v1/models",
        "openai": "https://api.openai.com/v1/models",
    }

    endpoint = endpoints.get(provider)
    if not endpoint:
        emit("error", code="UNKNOWN_PROVIDER", message=f"未知的模型提供商: {provider}")
        return

    try:
        resp = requests.get(
            endpoint,
            headers={"Authorization": f"Bearer {api_key}"},
            timeout=10,
        )
        if resp.status_code == 200:
            emit("result", status="success", message=f"{provider} API Key 验证成功")
        else:
            emit("error", code="INVALID_KEY",
                 message=f"API Key 无效: HTTP {resp.status_code}")
    except Exception as e:
        emit("error", code="CONNECTION_ERROR", message=f"连接失败: {str(e)}")


def handle_test_wechat(params):
    """测试微信公众号连接：获取 access_token 并返回结果"""
    import requests

    app_id = params.get("app_id", "")
    app_secret = params.get("app_secret", "")

    if not app_id or not app_secret:
        emit("error", code="MISSING_PARAMS", message="缺少 AppID 或 AppSecret")
        return

    # 获取当前出口 IP
    current_ip = ""
    try:
        ip_resp = requests.get("https://ifconfig.me/ip", timeout=5,
                               headers={"User-Agent": "curl/7.0"})
        if ip_resp.status_code == 200:
            current_ip = ip_resp.text.strip()
    except Exception:
        pass

    try:
        url = "https://api.weixin.qq.com/cgi-bin/token"
        resp = requests.get(url, params={
            "grant_type": "client_credential",
            "appid": app_id,
            "secret": app_secret,
        }, timeout=10)
        data = resp.json()
        # 脱敏日志：不记录 access_token
        safe_data = {k: ("***" if k == "access_token" else v)
                     for k, v in data.items()}
        logger.info("test_wechat response: %s",
                    json.dumps(safe_data, ensure_ascii=False)[:300])

        if data.get("access_token"):
            emit("result", status="success", ip=current_ip,
                 message="连接成功，access_token 获取正常")
        elif data.get("errcode") == 40164:
            # IP 不在白名单
            import re
            ip_match = re.search(r'invalid ip (\d+\.\d+\.\d+\.\d+)', data.get("errmsg", ""))
            real_ip = ip_match.group(1) if ip_match else current_ip or "未知"
            emit("result", status="ip_error", ip=real_ip,
                 message=f"当前出口 IP: {real_ip}（未在白名单中）")
        else:
            emit("error", code="WECHAT_ERROR",
                 message=f"错误 {data.get('errcode')}: {data.get('errmsg')}")
    except Exception as e:
        emit("error", code="CONNECTION_ERROR", message=f"连接失败: {str(e)}")


def handle_list_articles(params):
    """扫描 output 目录，列出已生成的文章"""
    output_dir = params.get("output_dir",
                            os.path.join(INK_HOME, "articles"))
    articles = []

    if os.path.exists(output_dir):
        for entry in sorted(os.listdir(output_dir), reverse=True):
            entry_path = os.path.join(output_dir, entry)
            # 查找目录中的 metadata JSON
            if os.path.isdir(entry_path):
                for fname in os.listdir(entry_path):
                    if fname.endswith("-metadata.json") or fname == "metadata.json":
                        meta_path = os.path.join(entry_path, fname)
                        try:
                            with open(meta_path, "r", encoding="utf-8") as f:
                                meta = json.load(f)
                                meta["id"] = entry
                                articles.append(meta)
                        except (json.JSONDecodeError, IOError):
                            pass  # 跳过损坏的 metadata 文件
                        break
            # 也支持直接放在 output_dir 下的 metadata 文件
            elif entry.endswith("-metadata.json"):
                meta_path = os.path.join(output_dir, entry)
                try:
                    with open(meta_path, "r", encoding="utf-8") as f:
                        meta = json.load(f)
                        meta["id"] = entry.replace("-metadata.json", "")
                        articles.append(meta)
                except (json.JSONDecodeError, IOError):
                    pass

    emit("result", status="success", articles=articles)


def handle_get_config(params):
    """读取 JSON 配置文件"""
    config_path = params.get("config_path",
                             os.path.join(INK_HOME, "config.json"))
    config = {}
    if os.path.exists(config_path):
        with open(config_path, "r", encoding="utf-8") as f:
            config = json.load(f)

    emit("result", status="success", config=config)


def handle_save_config(params):
    """保存 JSON 配置文件"""
    config_path = params.get("config_path",
                             os.path.join(INK_HOME, "config.json"))
    config_data = params.get("config", {})

    os.makedirs(os.path.dirname(config_path), exist_ok=True)
    with open(config_path, "w", encoding="utf-8") as f:
        json.dump(config_data, f, ensure_ascii=False, indent=2)

    emit("result", status="success", message="配置已保存")


def handle_read_file(params):
    """读取指定文件内容"""
    file_path = params.get("path", "")
    if not file_path or not os.path.exists(file_path):
        emit("error", code="FILE_NOT_FOUND", message=f"文件不存在: {file_path}")
        return
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()
        emit("result", status="success", content=content)
    except Exception as e:
        emit("error", code="READ_ERROR", message=str(e))


def handle_delete_article(params):
    """删除文章：根据 article_id 删除对应目录或文件"""
    import shutil

    article_id = params.get("article_id", "")
    output_dir = params.get("output_dir",
                            os.path.join(INK_HOME, "articles"))

    if not article_id:
        emit("error", code="MISSING_PARAMS", message="缺少 article_id")
        return

    deleted = []

    # 情况1：article_id 是子目录名
    dir_path = os.path.join(output_dir, article_id)
    if os.path.isdir(dir_path):
        shutil.rmtree(dir_path)
        deleted.append(dir_path)
    else:
        # 情况2：article_id 对应 output_dir 下的散落文件（以 id 为前缀）
        if os.path.exists(output_dir):
            for fname in os.listdir(output_dir):
                if fname.startswith(article_id):
                    fpath = os.path.join(output_dir, fname)
                    os.remove(fpath)
                    deleted.append(fpath)

    if deleted:
        emit("result", status="success",
             message=f"已删除 {len(deleted)} 个文件/目录", deleted=deleted)
    else:
        emit("error", code="NOT_FOUND",
             message=f"未找到文章: {article_id}")


def _upload_to_oss(local_path, oss_key, oss_config):
    """上传文件到阿里云 OSS"""
    import oss2
    auth = oss2.Auth(oss_config["oss_access_key_id"], oss_config["oss_access_key_secret"])
    endpoint = oss_config["oss_endpoint"]
    if not endpoint.startswith("http"):
        endpoint = f"https://{endpoint}"
    bucket = oss2.Bucket(auth, endpoint, oss_config["oss_bucket"])
    bucket.put_object_from_file(oss_key, local_path)


def handle_render_template(params):
    """用 AI 将纯文本渲染为微信公众号风格的 HTML 片段"""
    from llm_adapter import generate, LLMError

    text = params.get("text", "").strip()
    position = params.get("position", "footer")  # header or footer
    if not text:
        emit("error", code="MISSING_PARAMS", message="缺少 text 参数")
        return

    config = {}
    provider = params.get("provider", "deepseek")
    config["LLM_PROVIDER"] = provider
    for key in ["DEEPSEEK_API_KEY", "GLM_API_KEY", "DOUBAO_API_KEY",
                "KIMI_API_KEY", "OPENAI_API_KEY",
                "DEEPSEEK_MODEL", "GLM_MODEL", "DOUBAO_MODEL",
                "KIMI_MODEL", "OPENAI_MODEL"]:
        if params.get(key):
            config[key] = params[key]

    border = "border-bottom" if position == "header" else "border-top"
    prompt = f"""请将以下文字渲染为微信公众号文章{('头部' if position == 'header' else '尾部')}的 HTML 片段。

要求：
- 全部使用内联 CSS 样式
- 风格简洁美观，适合微信公众号
- 字体：-apple-system, BlinkMacSystemFont, 'PingFang SC', sans-serif
- 正文字号 14-15px，颜色 #333 或 #666
- 居中对齐
- 用 section 标签包裹，带 {border} 分隔线
- 只输出 HTML 代码，不要任何说明文字

文字内容：
{text}"""

    try:
        html = generate(prompt, config, timeout=60, need_search=False)
        if html:
            # 提取 HTML 部分
            from daily_ai_news import extract_html
            extracted = extract_html(html)
            emit("result", status="success", content=extracted or html)
        else:
            emit("error", code="RENDER_FAILED", message="渲染失败，未获得输出")
    except Exception as e:
        emit("error", code="RENDER_ERROR", message=str(e))


def handle_extract_files(params):
    """提取上传文件的文本内容，支持 Excel/PDF/Word"""
    file_paths = params.get("file_paths", [])
    if not file_paths:
        emit("error", code="MISSING_PARAMS", message="缺少 file_paths")
        return

    results = []
    for fpath in file_paths:
        if not os.path.exists(fpath):
            results.append({"path": fpath, "error": "文件不存在"})
            continue

        ext = os.path.splitext(fpath)[1].lower()
        name = os.path.basename(fpath)
        try:
            if ext in (".xlsx", ".xls"):
                text = _extract_excel(fpath)
            elif ext == ".pdf":
                text = _extract_pdf(fpath)
            elif ext in (".docx", ".doc"):
                text = _extract_docx(fpath)
            elif ext in (".pptx", ".ppt"):
                text = _extract_pptx(fpath)
            elif ext in (".txt", ".md", ".csv"):
                with open(fpath, "r", encoding="utf-8", errors="replace") as f:
                    text = f.read()
            else:
                results.append({"path": fpath, "name": name,
                                "error": f"不支持的文件类型: {ext}"})
                continue

            # 截断过长内容（避免超出 LLM 上下文）
            max_chars = 50000
            if len(text) > max_chars:
                text = text[:max_chars] + f"\n\n... (内容已截断，原文共 {len(text)} 字符)"

            results.append({"path": fpath, "name": name, "text": text,
                            "chars": len(text)})
            logger.info("Extracted %s: %d chars", name, len(text))
        except Exception as e:
            results.append({"path": fpath, "name": name, "error": str(e)})

    emit("result", status="success", files=results)


def _extract_excel(fpath):
    """提取 Excel 文件内容为文本表格"""
    import openpyxl
    wb = openpyxl.load_workbook(fpath, read_only=True, data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            rows.append(" | ".join(cells))
        if rows:
            parts.append(f"## Sheet: {sheet_name}\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(parts) if parts else "(空表格)"


def _extract_pdf(fpath):
    """提取 PDF 文件文本，pdfplumber 失败时回退到 PyMuPDF"""
    import pdfplumber
    parts = []
    with pdfplumber.open(fpath) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text and text.strip():
                parts.append(f"--- 第 {i+1} 页 ---\n{text}")
            # 提取表格
            tables = page.extract_tables()
            for t_idx, table in enumerate(tables):
                table_text = "\n".join(
                    " | ".join(str(c) if c else "" for c in row)
                    for row in table
                )
                parts.append(f"[表格 {t_idx+1}]\n{table_text}")

    if parts:
        logger.info("PDF extracted via pdfplumber: %d parts, %d chars",
                     len(parts), sum(len(p) for p in parts))
        return "\n\n".join(parts)

    # pdfplumber 提取失败，尝试 PyMuPDF 回退
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(fpath)
        for i, page in enumerate(doc):
            text = page.get_text()
            if text and text.strip():
                parts.append(f"--- 第 {i+1} 页 ---\n{text}")
        doc.close()
        if parts:
            logger.info("PDF extracted via PyMuPDF fallback: %d parts, %d chars",
                         len(parts), sum(len(p) for p in parts))
            return "\n\n".join(parts)
    except ImportError:
        logger.warning("PyMuPDF not installed, cannot fallback")
    except Exception as e:
        logger.warning("PyMuPDF fallback failed: %s", e)

    logger.warning("PDF extraction returned empty for: %s", fpath)
    return "(空 PDF — 无法提取文本，请尝试其他格式)"


def _extract_docx(fpath):
    """提取 Word 文档文本"""
    from docx import Document
    doc = Document(fpath)
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    # 提取表格
    for t_idx, table in enumerate(doc.tables):
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            parts.append(f"[表格 {t_idx+1}]\n" + "\n".join(rows))
    return "\n".join(parts) if parts else "(空文档)"


def _extract_pptx(fpath):
    """提取 PPT 文件内容为文本"""
    from pptx import Presentation
    prs = Presentation(fpath)
    parts = []
    for i, slide in enumerate(prs.slides):
        slide_parts = [f"--- 第 {i+1} 页 ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_parts.append(text)
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(" | ".join(cells))
                if rows:
                    slide_parts.append("[表格]\n" + "\n".join(rows))
        if len(slide_parts) > 1:
            parts.append("\n".join(slide_parts))
    return "\n\n".join(parts) if parts else "(空 PPT)"


def _cleanup_old_logs(max_days=7):
    """清理超过 max_days 天的日志文件"""
    import glob
    cutoff = datetime.now().timestamp() - max_days * 86400
    for f in glob.glob(os.path.join(LOG_DIR, "*.log")):
        if os.path.getmtime(f) < cutoff:
            try:
                os.remove(f)
            except OSError:
                pass


def _cleanup_old_cache(max_days=3):
    """清理超过 max_days 天的缓存文件"""
    import glob
    cutoff = datetime.now().timestamp() - max_days * 86400
    for f in glob.glob(os.path.join(CACHE_DIR, "*")):
        if os.path.getmtime(f) < cutoff:
            try:
                os.remove(f)
            except OSError:
                pass


def handle_get_logs(params):
    """读取最近的日志内容"""
    lines = params.get("lines", 100)
    log_file = os.path.join(LOG_DIR, f"{datetime.now().strftime('%Y-%m-%d')}.log")
    if not os.path.exists(log_file):
        emit("result", status="success", content="暂无日志", log_dir=LOG_DIR)
        return
    with open(log_file, "r", encoding="utf-8", errors="replace") as f:
        all_lines = f.readlines()
    tail = all_lines[-lines:]
    emit("result", status="success", content="".join(tail), log_dir=LOG_DIR)


def handle_clear_cache(params):
    """清理缓存目录"""
    import glob
    count = 0
    for f in glob.glob(os.path.join(CACHE_DIR, "*")):
        try:
            os.remove(f)
            count += 1
        except OSError:
            pass
    emit("result", status="success", message=f"已清理 {count} 个缓存文件")


def handle_publish_wechat(params):
    """发布文章到微信公众号草稿箱"""
    import requests as req
    from daily_ai_news import (
        get_access_token, upload_cover_image, create_draft,
    )
    from image_processor import process_images_in_html

    app_id = params.get("app_id", "")
    app_secret = params.get("app_secret", "")
    article_path = params.get("article_path", "")
    cover_path = params.get("cover_path", "")
    title = params.get("title", "未命名文章")
    author = params.get("author", "Ink")

    if not app_id or not app_secret:
        emit("error", code="MISSING_CONFIG", message="请先配置微信公众号 AppID 和 AppSecret")
        return
    if not article_path or not os.path.exists(article_path):
        emit("error", code="FILE_NOT_FOUND", message=f"文章文件不存在: {article_path}")
        return

    try:
        emit("progress", stage="publish", message="正在获取 access_token...")
        token = get_access_token(app_id, app_secret)

        emit("progress", stage="publish", message="正在读取文章...")
        html = open(article_path, "r", encoding="utf-8").read()

        # 上传文章内图片到微信 CDN
        emit("progress", stage="publish", message="正在上传文章图片到微信...")
        html = process_images_in_html(html, mode="wechat", access_token=token)

        # 上传封面图
        thumb_media_id = None
        if cover_path and os.path.exists(cover_path):
            emit("progress", stage="publish", message="正在上传封面图...")
            thumb_media_id = upload_cover_image(token, cover_path)

        if not thumb_media_id:
            logger.warning("封面图上传失败或未提供，微信草稿可能无法正常显示")
            emit("progress", stage="publish",
                 message="封面图缺失，尝试创建草稿...")

        # 创建草稿
        emit("progress", stage="publish", message="正在创建草稿...")
        media_id = create_draft(token, title, html, author, thumb_media_id)
        logger.info("create_draft returned media_id=%s", media_id)

        if media_id:
            # 更新 metadata 状态
            meta_dir = os.path.dirname(article_path)
            for f in os.listdir(meta_dir):
                if f.endswith("-metadata.json") or f == "metadata.json":
                    meta_path = os.path.join(meta_dir, f)
                    try:
                        meta = json.loads(open(meta_path, "r", encoding="utf-8").read())
                        meta["status"] = "published"
                        meta["wechat_media_id"] = media_id
                        open(meta_path, "w", encoding="utf-8").write(
                            json.dumps(meta, ensure_ascii=False, indent=2)
                        )
                    except Exception:
                        pass
                    break

            emit("result", status="success", message="已发布到草稿箱", media_id=media_id)
        else:
            emit("error", code="DRAFT_FAILED", message="创建草稿失败，请检查文章内容")

    except (SystemExit, RuntimeError):
        emit("error", code="AUTH_FAILED", message="access_token 获取失败，请检查 AppID/AppSecret 和 IP 白名单")
    except Exception as e:
        logger.exception("publish_wechat error")
        emit("error", code="PUBLISH_ERROR", message=str(e))


def _convert_html_to_format(html_content, target_ext, output_dir, timestamp=""):
    """将 HTML 内容转换为目标格式文件（PDF/DOCX/XLSX）。

    作为 Agent 未能生成原格式文件时的兜底方案。
    """
    import re
    os.makedirs(output_dir, exist_ok=True)

    # 从 HTML 提取纯文本（保留段落结构）
    def html_to_text(html):
        # 替换 <br> 为换行
        text = re.sub(r'<br\s*/?>', '\n', html)
        # 替换块级标签为换行
        text = re.sub(r'</(p|div|section|h[1-6]|li|tr)>', '\n', text)
        # 去掉所有 HTML 标签
        text = re.sub(r'<[^>]+>', '', text)
        # 解码 HTML 实体
        import html as html_mod
        text = html_mod.unescape(text)
        # 清理多余空行
        lines = [line.strip() for line in text.split('\n')]
        return '\n'.join(line for line in lines if line)

    text = html_to_text(html_content)
    if not text.strip():
        return None

    try:
        if target_ext == "pdf":
            return _text_to_pdf(text, output_dir, timestamp)
        elif target_ext in ("docx", "doc"):
            return _text_to_docx(text, output_dir, timestamp)
        elif target_ext in ("xlsx", "xls"):
            return _text_to_xlsx(text, output_dir, timestamp)
    except Exception as e:
        logger.error("Format conversion failed (%s): %s", target_ext, e)
    return None


def _text_to_pdf(text, output_dir, timestamp=""):
    """用 reportlab 将文本转为 PDF，支持中文。"""
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    import platform

    fname = f"translated-{timestamp}.pdf" if timestamp else "translated.pdf"
    out_path = os.path.join(output_dir, fname)

    # 注册中文字体
    font_candidates = []
    if platform.system() == "Darwin":
        font_candidates = [
            "/System/Library/Fonts/PingFang.ttc",
            "/System/Library/Fonts/STHeiti Light.ttc",
            "/Library/Fonts/Arial Unicode.ttf",
        ]
    elif platform.system() == "Windows":
        font_candidates = [
            "C:/Windows/Fonts/msyh.ttc",
            "C:/Windows/Fonts/simsun.ttc",
        ]
    else:
        font_candidates = [
            "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
            "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        ]

    font_name = "ChineseFont"
    registered = False
    for fp in font_candidates:
        if os.path.exists(fp):
            try:
                pdfmetrics.registerFont(TTFont(font_name, fp))
                registered = True
                break
            except Exception:
                continue

    used_font = font_name if registered else "Helvetica"
    c = canvas.Canvas(out_path, pagesize=A4)
    width, height = A4
    font_size = 11
    line_height = font_size * 1.8
    margin_x = 50
    margin_top = 50
    margin_bottom = 50
    max_width = width - 2 * margin_x
    y = height - margin_top

    c.setFont(used_font, font_size)

    for line in text.split("\n"):
        # 自动换行：按字符宽度拆分长行
        while line:
            # 估算一行能放多少字符（中文约 font_size 宽，英文约 font_size*0.6）
            chars_per_line = int(max_width / (font_size * 0.6))
            chunk = line[:chars_per_line]
            line = line[chars_per_line:]

            if y < margin_bottom:
                c.showPage()
                c.setFont(used_font, font_size)
                y = height - margin_top

            c.drawString(margin_x, y, chunk)
            y -= line_height

    c.save()
    return out_path


def _text_to_docx(text, output_dir, timestamp=""):
    """用 python-docx 将文本转为 Word 文档。"""
    from docx import Document
    fname = f"translated-{timestamp}.docx" if timestamp else "translated.docx"
    out_path = os.path.join(output_dir, fname)
    doc = Document()
    for line in text.split("\n"):
        if line.strip():
            doc.add_paragraph(line)
    doc.save(out_path)
    return out_path


def _text_to_xlsx(text, output_dir, timestamp=""):
    """用 openpyxl 将文本转为 Excel（每行一个单元格）。"""
    from openpyxl import Workbook
    fname = f"translated-{timestamp}.xlsx" if timestamp else "translated.xlsx"
    out_path = os.path.join(output_dir, fname)
    wb = Workbook()
    ws = wb.active
    ws.title = "Translation"
    for i, line in enumerate(text.split("\n"), 1):
        if line.strip():
            ws.cell(row=i, column=1, value=line)
    wb.save(out_path)
    return out_path


def main():
    # 启动时清理过期日志和缓存
    _cleanup_old_logs()
    _cleanup_old_cache()

    raw = sys.stdin.read()
    try:
        command = json.loads(raw)
    except json.JSONDecodeError:
        emit("error", code="INVALID_INPUT", message="无效的 JSON 输入")
        sys.exit(1)

    action = command.get("action")
    logger.info("action=%s", action)
    handlers = {
        "generate": handle_generate,
        "agent_generate": handle_agent_generate,
        "validate_key": handle_validate_key,
        "test_wechat": handle_test_wechat,
        "list_articles": handle_list_articles,
        "get_config": handle_get_config,
        "save_config": handle_save_config,
        "read_file": handle_read_file,
        "delete_article": handle_delete_article,
        "extract_files": handle_extract_files,
        "render_template": handle_render_template,
        "get_logs": handle_get_logs,
        "clear_cache": handle_clear_cache,
        "publish_wechat": handle_publish_wechat,
    }

    handler = handlers.get(action)
    if handler:
        try:
            handler(command)
        except Exception as e:
            logger.exception("handler %s failed", action)
            emit("error", code="INTERNAL_ERROR", message=str(e))
            sys.exit(1)
    else:
        logger.warning("unknown action: %s", action)
        emit("error", code="UNKNOWN_ACTION", message=f"未知操作: {action}")
        sys.exit(1)


if __name__ == "__main__":
    main()
