#!/usr/bin/env python3
"""
原格式保留翻译：在原文档中直接替换文字，保持格式和样式不变。

支持格式：
- Word (.docx) — 通过 python-docx 替换 run 文本
- PPT (.pptx) — 通过 python-pptx 替换 run 文本
- PDF (.pdf) — 通过 PyMuPDF 替换文本块
"""

import os
import logging
import json
import re

logger = logging.getLogger("ink")


def translate_file_inplace(input_path, output_path, target_lang, config,
                           emit_fn=None):
    """
    翻译文件，保持原格式。

    Args:
        input_path: 原始文件路径
        output_path: 输出文件路径
        target_lang: 目标语言描述，如 "英文"、"日文"
        config: LLM 配置 dict
        emit_fn: 进度回调 emit(type, **kwargs)

    Returns:
        output_path if success, None if failed
    """
    ext = os.path.splitext(input_path)[1].lower()

    if ext in (".docx", ".doc"):
        return _translate_docx(input_path, output_path, target_lang,
                               config, emit_fn)
    elif ext in (".pptx", ".ppt"):
        return _translate_pptx(input_path, output_path, target_lang,
                               config, emit_fn)
    elif ext == ".pdf":
        return _translate_pdf(input_path, output_path, target_lang,
                              config, emit_fn)
    else:
        logger.warning("Unsupported format for inplace translation: %s", ext)
        return None


# ============================================================
# LLM 批量翻译
# ============================================================

# 每批最大字符数（避免超出 LLM 上下文）
BATCH_MAX_CHARS = 8000
# 分隔符
SEP = "|||"


def _batch_translate(segments, target_lang, config, emit_fn=None):
    """
    批量翻译文本段落。

    将 segments 分批发送给 LLM，返回与 segments 等长的翻译列表。
    空字符串保持不变。
    """
    from llm_adapter import generate

    # 过滤出需要翻译的段落（非空且有实际文字内容）
    indexed = []  # (original_index, text)
    for i, seg in enumerate(segments):
        if seg and seg.strip():
            indexed.append((i, seg))

    if not indexed:
        return list(segments)

    # 分批
    batches = []
    current_batch = []
    current_chars = 0
    for idx, text in indexed:
        if current_chars + len(text) > BATCH_MAX_CHARS and current_batch:
            batches.append(current_batch)
            current_batch = []
            current_chars = 0
        current_batch.append((idx, text))
        current_chars += len(text)
    if current_batch:
        batches.append(current_batch)

    # 翻译结果
    results = list(segments)  # 复制原始列表
    total_batches = len(batches)

    for batch_idx, batch in enumerate(batches):
        if emit_fn:
            emit_fn("progress", stage="translating",
                    message=f"翻译中 ({batch_idx + 1}/{total_batches})...")

        texts = [text for _, text in batch]
        prompt = _build_translate_prompt(texts, target_lang)

        try:
            response = generate(prompt, config, timeout=120, need_search=False)
            translations = _parse_translate_response(response, len(texts))

            for i, (orig_idx, _) in enumerate(batch):
                if i < len(translations) and translations[i].strip():
                    results[orig_idx] = translations[i]
                # 翻译失败则保留原文

        except Exception as e:
            logger.error("Batch translation failed: %s", e)
            # 翻译失败保留原文

    return results


def _build_translate_prompt(texts, target_lang):
    """构建批量翻译 prompt"""
    numbered = "\n".join(f"[{i+1}] {t}" for i, t in enumerate(texts))
    return f"""你是专业翻译。请将以下编号文本翻译为{target_lang}。

要求：
1. 每条翻译用相同的编号格式输出：[编号] 翻译内容
2. 保持原文的语气和专业术语
3. 不要添加任何解释或注释
4. 数字、代码、URL 等保持原样
5. 如果原文只有标点或空白，保持原样

原文：
{numbered}

翻译："""


def _parse_translate_response(response, expected_count):
    """解析 LLM 翻译响应，提取编号翻译"""
    translations = [""] * expected_count

    # 匹配 [N] 翻译内容 格式（要求 [N] 在行首）
    for m in re.finditer(r'^\[(\d+)\]\s*(.+?)(?=^\[\d+\]|\Z)',
                         response, re.DOTALL | re.MULTILINE):
        idx = int(m.group(1)) - 1  # 转为 0-based
        text = m.group(2).strip()
        if 0 <= idx < expected_count:
            translations[idx] = text

    return translations


# ============================================================
# Word (.docx) 翻译
# ============================================================

def _translate_docx(input_path, output_path, target_lang, config, emit_fn):
    """翻译 Word 文档，保持原格式"""
    from docx import Document

    doc = Document(input_path)

    # 第一遍：收集所有文本段落
    segments = []
    segment_refs = []  # (type, ref) — 用于回写

    # 段落中的 runs
    for para in doc.paragraphs:
        for run in para.runs:
            if run.text and run.text.strip():
                segments.append(run.text)
                segment_refs.append(("run", run))

    # 表格中的 runs
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    for run in para.runs:
                        if run.text and run.text.strip():
                            segments.append(run.text)
                            segment_refs.append(("run", run))

    if not segments:
        logger.warning("No text found in docx: %s", input_path)
        return None

    if emit_fn:
        emit_fn("progress", stage="translating",
                message=f"Word 文档共 {len(segments)} 个文本段...")

    # 批量翻译
    translations = _batch_translate(segments, target_lang, config, emit_fn)

    # 回写翻译结果
    for i, (ref_type, ref) in enumerate(segment_refs):
        if ref_type == "run" and i < len(translations):
            ref.text = translations[i]

    doc.save(output_path)
    logger.info("Translated docx saved: %s (%d segments)", output_path,
                len(segments))
    return output_path


# ============================================================
# PPT (.pptx) 翻译
# ============================================================

def _translate_pptx(input_path, output_path, target_lang, config, emit_fn):
    """翻译 PPT，保持原格式"""
    from pptx import Presentation

    prs = Presentation(input_path)

    # 收集所有文本
    segments = []
    segment_refs = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text and run.text.strip():
                            segments.append(run.text)
                            segment_refs.append(("run", run))
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if run.text and run.text.strip():
                                    segments.append(run.text)
                                    segment_refs.append(("run", run))

    if not segments:
        logger.warning("No text found in pptx: %s", input_path)
        return None

    if emit_fn:
        emit_fn("progress", stage="translating",
                message=f"PPT 共 {len(segments)} 个文本段...")

    translations = _batch_translate(segments, target_lang, config, emit_fn)

    for i, (ref_type, ref) in enumerate(segment_refs):
        if ref_type == "run" and i < len(translations):
            ref.text = translations[i]

    prs.save(output_path)
    logger.info("Translated pptx saved: %s (%d segments)", output_path,
                len(segments))
    return output_path


# ============================================================
# PDF 翻译
# ============================================================

def _translate_pdf(input_path, output_path, target_lang, config, emit_fn):
    """翻译 PDF，尽量保持原格式。

    使用 PyMuPDF 逐页提取文本块，翻译后覆盖写回。
    注意：PDF 文本替换有局限性，复杂排版可能不完美。
    """
    try:
        import fitz
    except ImportError:
        logger.error("PyMuPDF not installed, cannot translate PDF inplace")
        return None

    doc = fitz.open(input_path)

    # 收集所有页面的文本块
    all_blocks = []  # [(page_idx, block_idx, text, rect, font_size)]
    segments = []

    for page_idx in range(len(doc)):
        page = doc[page_idx]
        blocks = page.get_text("dict")["blocks"]
        for b_idx, block in enumerate(blocks):
            if block["type"] != 0:  # 只处理文本块
                continue
            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    text = span.get("text", "").strip()
                    if text:
                        rect = fitz.Rect(span["bbox"])
                        font_size = span.get("size", 11)
                        font_name = span.get("font", "")
                        color = span.get("color", 0)
                        all_blocks.append((page_idx, rect, text,
                                           font_size, font_name, color))
                        segments.append(text)

    if not segments:
        logger.warning("No text found in PDF: %s", input_path)
        doc.close()
        return None

    if emit_fn:
        emit_fn("progress", stage="translating",
                message=f"PDF 共 {len(segments)} 个文本段...")

    translations = _batch_translate(segments, target_lang, config, emit_fn)

    # 回写：逐页处理
    if emit_fn:
        emit_fn("progress", stage="writing",
                message="正在写入翻译后的 PDF...")

    for i, (page_idx, rect, orig_text, font_size, font_name, color) in \
            enumerate(all_blocks):
        if i >= len(translations) or not translations[i].strip():
            continue
        if translations[i] == orig_text:
            continue

        page = doc[page_idx]

        # 用白色矩形覆盖原文
        page.add_redact_annot(rect)
        page.apply_redactions()

        # 选择字体：尝试使用中日韩字体
        use_font = "china-s"  # PyMuPDF 内置中文字体
        if any(ord(c) > 0x3000 for c in translations[i]):
            use_font = "china-s"

        # 插入翻译文本
        # 调整字号以适应原始区域
        adjusted_size = font_size
        if rect.height > 0:
            adjusted_size = min(font_size, max(rect.height * 0.8, 6.0))
        try:
            page.insert_textbox(
                rect, translations[i],
                fontsize=adjusted_size,
                fontname=use_font,
                color=_int_to_rgb(color),
                align=0,  # 左对齐
            )
        except Exception as e:
            logger.warning("Failed to insert text at page %d: %s",
                           page_idx, e)
            # 回退：用默认字体
            try:
                page.insert_textbox(
                    rect, translations[i],
                    fontsize=adjusted_size,
                    align=0,
                )
            except Exception:
                pass

    doc.save(output_path)
    doc.close()
    logger.info("Translated PDF saved: %s (%d segments)", output_path,
                len(segments))
    return output_path


def _int_to_rgb(color_int):
    """将整数颜色值转为 (r, g, b) 元组 (0-1 范围)"""
    if isinstance(color_int, (list, tuple)):
        return tuple(c / 255.0 for c in color_int[:3])
    r = ((color_int >> 16) & 0xFF) / 255.0
    g = ((color_int >> 8) & 0xFF) / 255.0
    b = (color_int & 0xFF) / 255.0
    return (r, g, b)
